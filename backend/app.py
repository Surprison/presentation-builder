import json
import ast
import os
from typing import Optional, List, Dict, Any
from datetime import datetime, timedelta

from fastapi import FastAPI, Depends, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware  # для работы CORS
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.responses import FileResponse
from pydantic import BaseModel
import openai
from dotenv import load_dotenv

from sqlalchemy import create_engine, Column, String, Integer, Text, ForeignKey
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship

from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from pptx import Presentation as PPTXPresentation
from pptx.util import Pt

from jose import JWTError, jwt
from passlib.context import CryptContext

# ----------------------------
# Загрузка переменных окружения и настройка OpenAI
# ----------------------------
load_dotenv(dotenv_path=".env")
api_key = os.getenv("OPENAI_API_KEY")
print(f"API Key: {api_key}")
openai.api_key = api_key

# ----------------------------
# Регистрация шрифта для PDF (корректный путь к TTF-файлу, например, arial.ttf)
# ----------------------------
pdfmetrics.registerFont(TTFont("DejaVuSans", "DejaVuSans.ttf"))

# Создадим собственные стили для заголовка и контента слайда, используя зарегистрированный шрифт
styles = getSampleStyleSheet() 

title_style = ParagraphStyle(
    'TitleStyle',
    parent=styles['Heading1'],
    fontName="DejaVuSans",
    fontSize=40,
    leading=55
)

content_style = ParagraphStyle(
    'ContentStyle',
    parent=styles['Normal'],
    fontName="DejaVuSans",
    fontSize=20,
    leading=35
)

# ----------------------------
# Настройка базы данных SQLite
# ----------------------------
DATABASE_URL = "sqlite:///./presentation.db"
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# Определение модели для презентаций
class PresentationModel(Base):
    __tablename__ = "presentations"
    id = Column(Integer, primary_key=True, index=True)
    theme = Column(String, index=True)
    slides = Column(Text)  # Хранение слайдов в формате JSON
    settings = Column(Text, nullable=True)  # Дополнительные настройки (JSON)
    num_pages = Column(Integer, nullable=True)  # Количество слайдов, если задано
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)  # Идентификатор владельца

    owner = relationship("UserModel", back_populates="presentations")

# Определение модели пользователя для авторизации
class UserModel(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True)
    hashed_password = Column(String)
    
    presentations = relationship("PresentationModel", back_populates="owner", cascade="all, delete")

Base.metadata.create_all(bind=engine)

# ----------------------------
# Создание экземпляра FastAPI
# ----------------------------
app = FastAPI()

# Настройка CORS (разрешаем все источники — для разработки)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # В продакшене укажите конкретные домены
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ----------------------------
# Функция для работы с БД (зависимость)
# ----------------------------
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ----------------------------
# Конфигурация авторизации (JWT)
# ----------------------------
SECRET_KEY = "your-secret-key"  # Замените на более надёжное значение!
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="users/login")

def get_password_hash(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return pwd_context.verify(plain_password, hashed_password)

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta or timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_user(db: Session, username: str) -> Optional[UserModel]:
    return db.query(UserModel).filter(UserModel.username == username).first()

# ----------------------------
# Зависимость для получения текущего пользователя
# ----------------------------
async def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials", 
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
    except JWTError:
        raise credentials_exception
    user = get_user(db, username=username)
    if user is None:
        raise credentials_exception
    return user

# ----------------------------
# Pydantic-модели для пользователей
# ----------------------------
class UserCreate(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class TokenData(BaseModel):
    username: Optional[str] = None

# ----------------------------
# Pydantic-модели для презентаций
# ----------------------------
class PresentationRequest(BaseModel):
    theme: str
    settings: Optional[Dict[str, Any]] = None
    num_pages: Optional[int] = None

class PresentationUpdate(BaseModel):
    theme: Optional[str] = None
    slides: Optional[List[Dict[str, Any]]] = None
    settings: Optional[Dict[str, Any]] = None
    num_pages: Optional[int] = None

# ----------------------------
# Эндпоинты для авторизации пользователей
# ----------------------------
@app.post("/users/register", response_model=Token)
def register(user: UserCreate, db: Session = Depends(get_db)):
    if get_user(db, user.username):
        raise HTTPException(status_code=400, detail="Username already registered")
    hashed_password = get_password_hash(user.password)
    new_user = UserModel(username=user.username, hashed_password=hashed_password)
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    access_token = create_access_token(data={"sub": new_user.username})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/users/login", response_model=Token)
def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = get_user(db, form_data.username)
    if not user or not verify_password(form_data.password, user.hashed_password):
        raise HTTPException(status_code=400, detail="Incorrect username or password")
    access_token = create_access_token(data={"sub": user.username})
    return {"access_token": access_token, "token_type": "bearer"}

@app.get("/users/me")
async def read_users_me(current_user: UserModel = Depends(get_current_user)):
    return {"username": current_user.username}

# ----------------------------
# Эндпоинты для работы с презентациями (доступны только владельцу)
# ----------------------------

# Создание презентации
@app.post("/generate")
async def generate_presentation(request: PresentationRequest, db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    if request.num_pages is not None:
        prompt = (
        f"Создай структуру презентации на тему \"{request.theme}\".\n"
        f"Презентация должна состоять из {request.num_pages} слайдов.\n"
        "Первый слайд должен быть титульным листом, содержащим только название темы как заголовок.\n"
        "Для каждого следующего слайда выведи ровно 5 строк:\n"
        "1. Заголовок\n"
        "2. Краткое описание (2-3 предложения)\n"
        "3. Основной пункт 1\n"
        "4. Основной пункт 2\n"
        "5. Основной пункт 3\n"
        "Только между слайдами должна быть пустая строка.\n"
        "ВАЖНО: Выводи только текст слайдов – без слов 'Слайд', 'Заголовок', 'Описание', 'Пункт', нумераций или других служебных меток.\n"
        "Начинай:"
    )
    else:
        prompt = (
        f"Создай структуру презентации на тему \"{request.theme}\".\n"
        "Первый слайд должен быть титульным листом, содержащим только название темы как заголовок.\n"
        "Для каждого слайда выведи ровно 5 строк:\n"
        "1. Заголовок\n"
        "2. Краткое описание (2-3 предложения)\n"
        "3. Основной пункт 1\n"
        "4. Основной пункт 2\n"
        "5. Основной пункт 3\n"
        "Только между слайдами должна быть пустая строка.\n"
        "ВАЖНО: Выводи исключительно текст слайдов без служебных меток (никаких слов типа 'Слайд', 'Заголовок', 'Описание', 'Пункт').\n"
        "Начинай:"
    )

    response = openai.ChatCompletion.create(
        model="gpt-4.1",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=1200
    )
    slides_text = response.choices[0].message.content.split("\n\n")
    slides = [
        {"title": slide.split("\n")[0], "content": slide.split("\n")[1:]}
        for slide in slides_text if slide
    ]

    # Сохраняем презентацию с привязкой к текущему пользователю
    presentation = PresentationModel(
        theme=request.theme,
        slides=json.dumps(slides),
        settings=json.dumps(request.settings) if request.settings is not None else None,
        num_pages=request.num_pages,
        user_id=current_user.id  # Привязываем презентацию к пользователю
    )
    db.add(presentation)
    db.commit()
    db.refresh(presentation)
    
    return {
        "id": presentation.id,
        "theme": presentation.theme,
        "slides": slides,
        "settings": request.settings,
        "num_pages": request.num_pages
    }

# Получить презентацию (только если текущий пользователь владелец)
@app.get("/presentation/{presentation_id}")
def get_presentation(presentation_id: int, db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    presentation = db.query(PresentationModel).filter(PresentationModel.id == presentation_id).first()
    if not presentation:
        raise HTTPException(status_code=404, detail="Презентация не найдена")
    if presentation.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Нет доступа к этой презентации")
    try:
        slides = json.loads(presentation.slides)
    except json.JSONDecodeError:
        slides = ast.literal_eval(presentation.slides)
    
    settings = None
    if presentation.settings:
        try:
            settings = json.loads(presentation.settings)
        except json.JSONDecodeError:
            settings = ast.literal_eval(presentation.settings)
    
    return {
        "id": presentation.id,
        "theme": presentation.theme,
        "slides": slides,
        "settings": settings,
        "num_pages": presentation.num_pages
    }

# Обновление презентации (только если текущий пользователь владелец)
@app.put("/presentation/{presentation_id}")
def update_presentation(presentation_id: int, update: PresentationUpdate, db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    presentation = db.query(PresentationModel).filter(PresentationModel.id == presentation_id).first()
    if not presentation:
        raise HTTPException(status_code=404, detail="Презентация не найдена")
    if presentation.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Нет доступа к этой презентации")
    
    if update.theme is not None:
        presentation.theme = update.theme
    if update.slides is not None:
        presentation.slides = json.dumps(update.slides)
    if update.settings is not None:
        presentation.settings = json.dumps(update.settings)
    if update.num_pages is not None:
        presentation.num_pages = update.num_pages
    
    db.commit()
    db.refresh(presentation)
    
    try:
        slides = json.loads(presentation.slides)
    except json.JSONDecodeError:
        slides = ast.literal_eval(presentation.slides)
    
    settings = None
    if presentation.settings:
        try:
            settings = json.loads(presentation.settings)
        except json.JSONDecodeError:
            settings = ast.literal_eval(presentation.settings)
    
    return {
        "message": "Презентация обновлена",
        "presentation": {
            "id": presentation.id,
            "theme": presentation.theme,
            "slides": slides,
            "settings": settings,
            "num_pages": presentation.num_pages
        }
    }

# Удаление презентации (только если текущий пользователь владелец)
@app.delete("/presentation/{presentation_id}")
def delete_presentation(presentation_id: int, db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    presentation = db.query(PresentationModel).filter(PresentationModel.id == presentation_id).first()
    if not presentation:
        raise HTTPException(status_code=404, detail="Презентация не найдена")
    if presentation.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Нет доступа к этой презентации")
    db.delete(presentation)
    db.commit()
    return {"message": "Презентация удалена"}

# Получить список презентаций текущего пользователя
@app.get("/presentations")
def list_presentations(db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    presentations = db.query(PresentationModel).filter(PresentationModel.user_id == current_user.id).order_by(PresentationModel.id.desc()).all()
    result = []
    for p in presentations:
        settings = None
        if p.settings:
            try:
                settings = json.loads(p.settings)
            except json.JSONDecodeError:
                settings = ast.literal_eval(p.settings)
        result.append({
            "id": p.id,
            "theme": p.theme,
            "num_pages": p.num_pages,
            "settings": settings
        })
    return result


@app.get("/export/pdf/{presentation_id}")
def export_pdf(presentation_id: int, 
               db: Session = Depends(get_db), 
               current_user: UserModel = Depends(get_current_user)):
    # Находим презентацию в базе данных
    presentation = db.query(PresentationModel).filter(PresentationModel.id == presentation_id).first()
    if not presentation:
        raise HTTPException(status_code=404, detail="Презентация не найдена")
    if presentation.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Нет доступа к этой презентации")

    # Пробуем распарсить слайды из поля presentation.slides
    try:
        slides = json.loads(presentation.slides)
    except json.JSONDecodeError:
        slides = ast.literal_eval(presentation.slides)

    export_dir = "./exports"
    os.makedirs(export_dir, exist_ok=True)
    pdf_filename = f"presentation_{presentation_id}.pdf"
    pdf_path = os.path.join(export_dir, pdf_filename)

    # Создаем объект canvas с альбомной ориентацией для страницы A4
    c = canvas.Canvas(pdf_path, pagesize=landscape(A4))
    width, height = landscape(A4)
    
    # Задаем координаты отступа и максимально допустимую ширину текста
    x_pos = 100
    y_pos = height - 100
    max_width = width * 0.8

    # Обрабатываем каждый слайд презентации
    for slide in slides:
        # Отрисовка заголовка слайда с помощью Paragraph
        p_title = Paragraph(slide["title"], title_style)
        w, h = p_title.wrap(max_width, height)
        p_title.drawOn(c, x_pos, y_pos - h)
        y_pos -= h + 10

        # Отрисовка содержимого слайда: для каждой строки создается Paragraph,
        # который автоматически переносит текст по заданной ширине
        for line in slide["content"]:
            p_content = Paragraph(line, content_style)
            w, h = p_content.wrap(max_width, height)
            # Если текст уходит за нижнюю границу, начинаем новую страницу
            if y_pos - h < 50:
                c.showPage()
                y_pos = height - 100
            p_content.drawOn(c, x_pos, y_pos - h)
            y_pos -= h + 4

        c.showPage()
        y_pos = height - 100
    c.save()

    return FileResponse(
        path=pdf_path,
        media_type="application/pdf",
        filename=pdf_filename
    )
    return {"message": "PDF создан", "file": pdf_path}

# Экспорт презентации в PPTX (проверка доступа)
@app.get("/export/pptx/{presentation_id}")
def export_pptx(presentation_id: int, db: Session = Depends(get_db), current_user: UserModel = Depends(get_current_user)):
    presentation = db.query(PresentationModel).filter(PresentationModel.id == presentation_id).first()
    if not presentation:
        raise HTTPException(status_code=404, detail="Презентация не найдена")
    if presentation.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Нет доступа к этой презентации")
    try:
        slides = json.loads(presentation.slides)
    except json.JSONDecodeError:
        slides = ast.literal_eval(presentation.slides)
    
    export_dir = "./exports"
    os.makedirs(export_dir, exist_ok=True)
    ppt_filename = f"presentation_{presentation_id}.pptx"
    ppt_path = os.path.join(export_dir, ppt_filename)
    
    
    prs = PPTXPresentation()
    for slide in slides:
        slide_layout = prs.slide_layouts[1]  # Layout: Title and Content
        ppt_slide = prs.slides.add_slide(slide_layout)
        title = ppt_slide.shapes.title
        content_placeholder = ppt_slide.placeholders[1]
        title.text = slide["title"]
        content_placeholder.text = "\n".join(slide["content"])
        for paragraph in content_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(25)
    
    # Сохраняем презентацию
    prs.save(ppt_path)
    
    # Возвращаем файл с корректными HTTP‑заголовками
    return FileResponse(
        path=ppt_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=ppt_filename
    )

    return {"message": "PPTX создан", "file": ppt_path}